import os

import pptx
from pptx import Presentation
import datetime

from src.data import load_data
from src.utils import change_pic, clear_text, change_text


def run(ppt, data):
    df = load_data(data)
    print(df)
    prs = Presentation(ppt)

    slides = prs.slides

    referencia = "Entrega"
    # Rename Slides
    for slide in slides:
        for shape in slide.shapes:
            if shape.name == referencia:
                slide.name = shape.text

    sl = 0
    for slide in slides:
        frame = 0
        for shape in slide.shapes:
            # Se não o shape não for picture, troca o texto
            if not type(shape) == pptx.shapes.picture.Picture:
                if shape.name in df.columns:
                    clear_text(shape)
                    if str(df.loc[df[referencia] == slide.name][shape.name].iloc[0]) == "nan":
                        change_text(shape, "")
                    else:
                        change_text(shape, str(df.loc[df[referencia] == slide.name][shape.name].iloc[0]).replace(';',
                                                                                                                 "").strip())
                    frame += 1

            if type(shape) == pptx.shapes.picture.Picture:
                if shape.name in df.columns:
                    change_pic(slide.shapes, shape, df.loc[df[referencia] == slide.name][shape.name].iloc[0])

            if shape.name == "HomeButton":
                shape.click_action.target_slide = slides[df.loc[df[referencia] == slide.name]['hyperlink'].iloc[0]]
        sl += 1

    hoje = datetime.date.today().strftime("%d-%m-%Y")

    filename = f'Generated-{hoje}.pptx'
    prs.save(filename)
    os.system(filename)
